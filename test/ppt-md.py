#!/usr/bin/env python3
import argparse
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def safe_filename(text):
    """Convert text to safe filename."""
    return re.sub(r'[^\w\s-]', '', text).strip()

def extract_slide_content(slide, image_dir=None, slide_number=1):
    """Extract text and images from a slide."""
    content = []
    
    # Handle slide title
    if slide.shapes.title:
        content.append(f"## {slide.shapes.title.text}\n")
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    content.append(f"{paragraph.text}\n")
                    
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_dir:
            try:
                # Save image and create markdown link
                image_path = os.path.join(
                    image_dir, 
                    f"slide_{slide_number}_{len(content)}.png"
                )
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)
                content.append(f"![Slide {slide_number} Image]({image_path})\n")
            except Exception as e:
                content.append(f"*[Image extraction failed: {str(e)}]*\n")
    
    return "\n".join(content)

def ppt_to_markdown(ppt_path, output_dir=None):
    """Convert PowerPoint to Markdown."""
    try:
        prs = Presentation(ppt_path)
        
        # Setup output directory
        if output_dir is None:
            output_dir = os.path.dirname(os.path.abspath(ppt_path))
        
        base_name = os.path.splitext(os.path.basename(ppt_path))[0]
        md_path = os.path.join(output_dir, f"{base_name}.md")
        image_dir = os.path.join(output_dir, f"{base_name}_images")
        
        # Create image directory if needed
        os.makedirs(image_dir, exist_ok=True)
        
        # Process slides
        markdown_content = [f"# {base_name}\n"]
        for i, slide in enumerate(prs.slides, 1):
            markdown_content.append(
                extract_slide_content(slide, image_dir, i)
            )
        
        # Write markdown file
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(markdown_content))
        
        return md_path
        
    except Exception as e:
        print(f"Error converting PowerPoint: {str(e)}")
        return None

def main():
    parser = argparse.ArgumentParser(
        description='Convert PowerPoint files to Markdown'
    )
    parser.add_argument('file', help='Path to PowerPoint file (.ppt or .pptx)')
    parser.add_argument(
        '--output-dir', '-o',
        help='Output directory (optional)',
        default=None
    )
    
    args = parser.parse_args()
    
    if not args.file.lower().endswith(('.ppt', '.pptx')):
        print("Error: File must be a PowerPoint file (.ppt or .pptx)")
        return
    
    output_path = ppt_to_markdown(args.file, args.output_dir)
    if output_path:
        print(f"Successfully converted to: {output_path}")

if __name__ == "__main__":
    main()
